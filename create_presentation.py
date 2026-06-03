#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

DARK_BLUE = RGBColor(30, 58, 138)
LIGHT_BLUE = RGBColor(66, 133, 244)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(64, 64, 64)
ACCENT = RGBColor(251, 188, 4)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, bullets_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE

    title_frame = header.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)

def add_screenshot_slide(title, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = LIGHT_BLUE
    header.line.color.rgb = LIGHT_BLUE

    title_frame = header.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    rect = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(5.3))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(240, 240, 240)
    rect.line.color.rgb = RGBColor(200, 200, 200)
    rect.line.width = Pt(2)

    text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CAPTURA AQUI\n\n" + caption
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

# PORTADA
add_title_slide("RETO JUEGO BRYAN", "Sistema de Combate y Disparo\nETAPA 2 - Unity 3D")

# PARTE 1
add_content_slide("PARTE 1: CONFIGURACION Y ARQUITECTURA", [
    "Expositor 1",
    "Estructura tecnica del proyecto",
    "Organizacion de assets y scripts",
    "Arquitectura de componentes"
])

add_content_slide("1.1 Organizacion de Carpetas", [
    "Assets/",
    "  + Models/ -> FBX Player, Cleric",
    "  + Texture/ -> Texturas PBR",
    "  + Scripts/ -> Codigo C#",
    "  + Editor/ -> Herramientas",
    "  + Animation/ -> Animator Controller"
])

add_screenshot_slide("1.1 Project Window", "Mostrar estructura de carpetas expandidas")

add_content_slide("1.2 Configuracion del Player", [
    "CharacterController -> Capsula colisiones",
    "Animator -> Animaciones (Blend Tree)",
    "Tag: 'Player' -> Identificacion",
    "Visual (hijo) -> Modelo 3D separado",
    "Cinemachine FreeLook -> Camara orbital"
])

add_screenshot_slide("1.2 Player Inspector", "Jerarquia del Player con componentes")

add_content_slide("1.3 Segundo Personaje (TargetDummy)", [
    "CapsuleCollider -> Detecta impactos",
    "  - height: 2m | radius: 0.5m",
    "TargetDummy.cs -> IDamageable",
    "Tag: 'Target' -> Identificacion objetivo",
    "Auto-respawn -> Cada 3 segundos"
])

add_screenshot_slide("1.3 TargetDummy Setup", "Cleric en escena con CapsuleCollider")

add_content_slide("1.4 Tags vs Layers", [
    "TAGS usados:",
    "  + 'Player' -> Ignora sus propias balas",
    "  + 'Target' -> Identifica objetivos",
    "",
    "Por que Tags?",
    "  + Suficiente para esta escala",
    "  + Mas simple que Layers"
])

add_content_slide("1.5 Arquitectura de Scripts", [
    "Scripts/Player/ -> Movimiento y animacion",
    "Scripts/Combat/ -> Sistema dano compartido",
    "Scripts/Enemy/ -> Salud de enemigos",
    "Scripts/UI/ -> HUD e interfaz",
    "Editor/ -> Automatizacion editor"
])

add_screenshot_slide("1.5 Estructura Scripts", "Carpetas expandidas en Project window")

add_content_slide("1.6 Responsabilidad de Scripts", [
    "PlayerMovement -> WASD + Salto",
    "PlayerAnimator -> Blend Tree (Speed, Jump)",
    "WeaponSystem -> Disparo + Municion",
    "Projectile -> Fisica + Impacto",
    "IDamageable -> Interfaz dano",
    "HUDDisplay -> UI tiempo real"
])

# PARTE 2
add_content_slide("PARTE 2: SISTEMA DE DISPARO Y COMBATE", [
    "Expositor 2",
    "Flujo del combate",
    "Fisica de proyectiles",
    "Efectos visuales procedurales"
])

add_content_slide("2.1 Comunicacion entre Scripts", [
    "Patron Observer (Eventos C#):",
    "  1. WeaponSystem dispara evento",
    "  2. HUD se suscribe -> se actualiza",
    "  3. ArmAnimation se suscribe -> anima",
    "",
    "Desacoplamiento total",
    "No necesitan conocerse"
])

add_screenshot_slide("2.1 Diagrama Flujo", "Disparo -> OnShoot -> HUD/Animation")

add_content_slide("2.2 WeaponSystem: Shoot()", [
    "Flujo:",
    "  1. Verificar municion (> 0)",
    "  2. Verificar fire rate",
    "  3. CurrentAmmo--",
    "  4. OnAmmoChanged.Invoke()",
    "  5. BulletFactory.Create()",
    "  6. Aplicar velocidad"
])

add_screenshot_slide("2.2 Shoot Code", "Codigo en Visual Studio destacando flujo")

add_content_slide("2.3 Municion y Recarga", [
    "Magazine: 6 balas",
    "Fire Rate: 0.5 segundos",
    "",
    "Recarga con R:",
    "  + Solo si CurrentAmmo < 6",
    "  + Vuelve a 6 instantaneamente",
    "  + No dispara durante recarga"
])

add_screenshot_slide("2.3 HUD Municion", "BALAS: 3 / 6 y RECARGANDO...")

add_content_slide("2.4 Disparo 100% Horizontal", [
    "Problema: disparaba hacia arriba",
    "",
    "Solucion:",
    "  + dir = transform.forward",
    "  + dir.y = 0 -> fuerza horizontal",
    "  + dir.Normalize() -> re-normaliza",
    "  + rb.useGravity = false"
])

add_screenshot_slide("2.4 Gizmo Disparo", "Green line mostrando direccion 5m")

add_content_slide("2.5 Sistema de Espada", [
    "FindHandBone() -> busca RightHand",
    "SetParent(handBone, false) -> local",
    "localPosition = (0, 0.05, 0) -> ajuste",
    "",
    "Espada sigue animacion automaticamente",
    "Sin codigo extra de movimiento"
])

add_screenshot_slide("2.5 Warrior con Espada", "Personaje enganchando arma correctamente")

add_content_slide("2.6 Anatomia del Proyectil", [
    "[RequireComponent] Rigidbody + Collider",
    "",
    "OnCollisionEnter():",
    "  + ContactPoint = punto exacto",
    "  + collision.normal = orientacion",
    "",
    "OnTriggerEnter(): Fallback"
])

add_screenshot_slide("2.6 Proyectil Physics", "Esfera amarilla impactando TargetDummy")

add_content_slide("2.7 Sistema Dano (IDamageable)", [
    "Interface publica:",
    "  + TakeDamage(float amount)",
    "  + Die()",
    "",
    "Implementan: PlayerHealth, Enemy, Target",
    "Polimorfismo: cualquier recibe dano"
])

add_screenshot_slide("2.7 IDamageable Code", "Interface y implementaciones")

add_content_slide("2.8 Efectos Visuales (VFX)", [
    "ParticleSystem procedural (sin Prefab):",
    "  + 25 particulas en burst",
    "  + Emision conica segun superficie",
    "  + Color orange -> red fade",
    "  + Luz puntual 3 intensidad",
    "  + Destruye en 1 segundo"
])

add_screenshot_slide("2.8 VFX Impacto", "Particulas naranjas explosion radial")

# PARTE 3
add_content_slide("PARTE 3: UI, PRACTICAS Y CONCLUSIONES", [
    "Expositor 3",
    "Interfaz de usuario",
    "Arquitectura profesional",
    "Principios SOLID aplicados"
])

add_content_slide("3.1 Sistema UI (OnGUI)", [
    "Immediate Mode GUI (IMGUI):",
    "  + Rapido: sin Canvas",
    "  + Procedural: todo codigo",
    "",
    "HUDDisplay.cs:",
    "  + Dibuja rectangulos",
    "  + Se suscribe eventos"
])

add_screenshot_slide("3.1 HUD Gameplay", "BALAS, vida, mira, mensajes")

add_content_slide("3.2 Elementos del HUD", [
    "Municion: BALAS: X/6 (rojo si vacio)",
    "Barra vida: coloreada (rojo->verde)",
    "Crosshair: cruz centro pantalla",
    "Mensajes: SIN BALAS - Presiona R",
    "",
    "Todos actualizados por eventos"
])

add_content_slide("3.3 Buenas Practicas", [
    "Interfaces (IDamageable)",
    "Eventos (Observer Pattern)",
    "Factory Method (BulletFactory)",
    "Guard Clauses (early returns)",
    "Composicion sobre herencia",
    "Separacion responsabilidades"
])

add_content_slide("3.4 Conceptos Clave Unity", [
    "MonoBehaviour Lifecycle:",
    "  Awake() -> Start() -> Update() -> FixedUpdate()",
    "",
    "Fisica: Rigidbody + Colliders",
    "Espacios: localPosition vs position",
    "Vectores: transform.forward, normalize()"
])

add_content_slide("3.5 Patrones Programacion", [
    "Observer -> Eventos para desacoplamiento",
    "Factory Method -> BulletFactory.Create()",
    "Polimorfismo -> IDamageable.TakeDamage()",
    "Guard Clauses -> Early returns",
    "Singleton -> Camera.main"
])

add_content_slide("3.6 Cumplimiento Rubrica", [
    "Arquitectura: Codigo limpio, eventos",
    "Mecanicas: Movimiento fluido, sin bugs",
    "Multimedia: Modelos, animaciones, VFX",
    "Optimizacion: 60 FPS, cleanup",
    "POO: Interfaces, encapsulacion"
])

add_screenshot_slide("3.6 Rubrica Completada", "Evidencia en-game todos criterios")

add_content_slide("3.7 Principios SOLID", [
    "S - Single Responsibility: cada script una funcion",
    "O - Open/Closed: extensible interfaces",
    "L - Liskov Substitution: cualquier IDamageable",
    "I - Interface Segregation: solo necesario",
    "D - Dependency Inversion: eventos"
])

add_content_slide("CONCLUSION", [
    "Arquitectura profesional desde inicio",
    "Desacoplamiento total componentes",
    "Codigo escalable y mantenible",
    "Implementacion patrones reales",
    "Rubrica 100% cumplida",
    "",
    "Listo para produccion."
])

prs.save('RETO_JUEGO_PRESENTACION.pptx')
print("Presentacion creada: RETO_JUEGO_PRESENTACION.pptx")
print("Total slides:", len(prs.slides))
